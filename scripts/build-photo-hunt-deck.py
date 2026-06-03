#!/usr/bin/env python3
"""Build Photo Hunt marketing slide deck with embedded screenshots."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "marketing" / "screenshots"
OUT = ROOT / "docs" / "marketing" / "Photo-Hunt-Overview.pptx"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SLATE = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_block(slide, title: str, subtitle: str = "", top=0.45):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(8)


def add_bullets(slide, items: list[str], left=0.75, top=1.55, width=5.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)


def add_image(slide, filename: str, left, top, width, height=None):
    path = SHOTS / filename
    if not path.exists():
        note = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        note.text_frame.text = f"[Screenshot missing: {filename}]"
        return
    if height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def add_footer(slide, text="familyphotohunt.com"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = SLATE


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — Title
    s = blank_slide(prs)
    hero = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(3))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Photo Hunt"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Turn any gathering into a themed photo contest"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(191, 219, 254)
    p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = "www.familyphotohunt.com"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_before = Pt(24)

    # 2 — What it is
    s = blank_slide(prs)
    add_title_block(s, "What is Photo Hunt?", "A web app for group photo contests — no download required.")
    add_bullets(
        s,
        [
            "Organizer creates a contest with themed categories (Food, Best View, Silly Moment…).",
            "Participants join with a 4-character code from their phone browser.",
            "Everyone uploads photos, votes on favorites, and celebrates winners together.",
            "Perfect for family trips, reunions, retreats, and team events.",
        ],
        width=11.8,
    )
    add_footer(s)

    # 3 — Four stages
    s = blank_slide(prs)
    add_title_block(s, "Four stages, one smooth flow", "The organizer controls timing; participants always know what to do.")
    stages = [
        ("1  Setup", "Add categories. Join code hidden until ready."),
        ("2  Photo Collection", "Share join code. Participants upload & submit photos."),
        ("3  Voting", "Collection closes. One vote per category (not your own)."),
        ("4  Results", "Winners, scoreboard, TV reveal, and downloads."),
    ]
    y = 1.65
    for label, desc in stages:
        box = s.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(0.95))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = NAVY
        y += 1.15
    add_footer(s)

    # 4 — Home / join
    s = blank_slide(prs)
    add_title_block(s, "Joining is simple", "Enter a code on the home page — or follow a link from your organizer.")
    add_bullets(
        s,
        [
            "4-character join codes — easy to text to a group.",
            "Create an account once; stay signed in on the same device.",
            "Contest cards show location, date, and current stage.",
        ],
        width=5.5,
    )
    add_image(s, "01-home-desktop.png", 6.5, 1.45, 6.2)
    add_footer(s)

    # 5 — Login
    s = blank_slide(prs)
    add_title_block(s, "Sign in or register in seconds", "Name, email, and password — then you're in the contest.")
    add_image(s, "02-login-desktop.png", 0.75, 1.45, 5.9)
    add_bullets(
        s,
        [
            "Works on phone and desktop browsers.",
            "Join code can be pre-filled from an invite link.",
            "Password show/hide for easier mobile entry.",
        ],
        left=7.0,
        width=5.5,
    )
    add_footer(s)

    # 6 — Setup
    s = blank_slide(prs)
    add_title_block(s, "Stage 1 — Setup", "Build your category list before anyone joins.")
    add_bullets(
        s,
        [
            "Create contest with location/event name + month/year.",
            "Add suggested categories or custom themes with descriptions.",
            "Advance to Photo Collection when categories are ready.",
            "Join code appears only after leaving Setup.",
        ],
        width=5.6,
    )
    add_image(s, "04-help-admin-desktop.png", 6.4, 1.35, 6.3)
    add_footer(s)

    # 7 — Collection
    s = blank_slide(prs)
    add_title_block(s, "Stage 2 — Open Photo Collection", "Invite the group and collect photos by category.")
    add_bullets(
        s,
        [
            "Copy a ready-made invite message (link + join code).",
            "Participants upload from camera or photo library.",
            "Rank multiple photos — top photo is submitted per category.",
            "Organizer tracks who submitted and category progress.",
        ],
        width=11.8,
    )
    add_footer(s)

    # 8 — Voting
    s = blank_slide(prs)
    add_title_block(s, "Stage 3 — Voting", "Fair, simple, one favorite per category.")
    add_bullets(
        s,
        [
            "Gallery view per category — browse all submitted photos.",
            "One vote per person per category.",
            "Cannot vote for your own photo.",
            "Organizer sees voting progress bar and per-participant status.",
            "Ties for most votes can share a win.",
        ],
        width=11.8,
    )
    add_footer(s)

    # 9 — Results
    s = blank_slide(prs)
    add_title_block(s, "Stage 4 — Results", "Celebrate winners and keep the memories.")
    add_bullets(
        s,
        [
            "Full scoreboard with vote counts per category.",
            "Start Winner Reveal — TV-friendly slideshow for the room.",
            "Download all photos or winners-only ZIP.",
            "Participants: Save to Photos on mobile from results.",
        ],
        width=11.8,
    )
    add_footer(s)

    # 10 — Participant guide
    s = blank_slide(prs)
    add_title_block(s, "Participant experience", "Built-in guide walks users through every step.")
    add_image(s, "05-help-participants-desktop.png", 0.65, 1.35, 6.4)
    add_bullets(
        s,
        [
            "Join → Upload → Submit → Vote → Results",
            "Skip categories you don't have photos for.",
            "Ready to Submit locks uploads — plan before confirming.",
            "Help at /help/participants anytime.",
        ],
        left=7.2,
        width=5.4,
    )
    add_footer(s)

    # 11 — Organizer tools
    s = blank_slide(prs)
    add_title_block(s, "Organizer tools", "Everything you need to run the contest from your phone or laptop.")
    add_bullets(
        s,
        [
            "Visual stage stepper — move forward or back with confirmation.",
            "Copyable announcement snippets for each stage.",
            "Participant panel with submission & voting status.",
            "Suggested category library (Animal, Architecture, Food…).",
            "Admin dashboard lists all your contests with current stage.",
        ],
        width=11.8,
    )
    add_footer(s)

    # 12 — Mobile
    s = blank_slide(prs)
    add_title_block(s, "Built for phones", "Upload, vote, and save winners from the mobile browser.")
    add_image(s, "01-home-mobile.png", 0.8, 1.35, 3.2)
    add_image(s, "02-login-mobile.png", 4.3, 1.35, 3.2)
    add_bullets(
        s,
        [
            "Touch-friendly buttons and galleries.",
            "Camera roll upload — no separate app.",
            "Save winning photos to camera roll on iOS/Android.",
            "Desktop great for admin + projector winner reveal.",
        ],
        left=7.8,
        top=1.8,
        width=4.8,
    )
    add_footer(s)

    # 13 — Feature table
    s = blank_slide(prs)
    add_title_block(s, "Feature highlights")
    rows = [
        ("Join codes", "4 characters — text to your group"),
        ("Categories", "Custom themes + optional descriptions"),
        ("Submissions", "One photo per person per category"),
        ("Voting", "One vote per category; no self-voting"),
        ("Reveal", "Fullscreen winner presentation"),
        ("Downloads", "ZIP all photos or winners only"),
    ]
    y = 1.55
    for label, detail in rows:
        box = s.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(15)
        p2.font.color.rgb = NAVY
        y += 0.85
    add_footer(s)

    # 14 — Getting started
    s = blank_slide(prs)
    add_title_block(s, "Get started", "Free to try — create your first contest in minutes.")
    add_bullets(
        s,
        [
            "1.  Visit www.familyphotohunt.com",
            "2.  Sign in and tap Create or manage a contest",
            "3.  Add categories and invite your group with the join code",
            "4.  Move through stages as your event unfolds",
        ],
        width=11.8,
        size=22,
    )
    cta = s.shapes.add_textbox(Inches(0.75), Inches(5.0), Inches(11.8), Inches(0.8))
    p = cta.text_frame.paragraphs[0]
    p.text = "Questions? Share this deck and the Photo-Hunt-Prospect-Guide.md companion doc."
    p.font.size = Pt(16)
    p.font.color.rgb = SLATE
    add_footer(s)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
